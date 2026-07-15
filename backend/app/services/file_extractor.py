from pypdf import PdfReader
from docx import Document
from pptx import Presentation



def extract_text(file, filename):

    filename = filename.lower()



    # TXT FILE
    if filename.endswith(".txt"):

        content = file.read()

        return content.decode("utf-8")




    # PDF FILE
    elif filename.endswith(".pdf"):

        reader = PdfReader(file)

        text = ""

        for page in reader.pages:

            text += page.extract_text() or ""

        return text





    # DOCX FILE
    elif filename.endswith(".docx"):

        doc = Document(file)

        text = ""

        for paragraph in doc.paragraphs:

            text += paragraph.text + "\n"

        return text





    # PPTX FILE
    elif filename.endswith(".pptx"):

        presentation = Presentation(file)

        text = ""

        for slide in presentation.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text += shape.text + "\n"

        return text





    else:

        raise ValueError(
            "Unsupported file type"
        )